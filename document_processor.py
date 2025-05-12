from typing import List, Union
import PyPDF2
from pptx import Presentation
import os

class DocumentProcessor:
    def __init__(self):
        self.supported_extensions = ['.pdf', '.pptx', '.ppt']

    def process_file(self, file_path: str) -> List[str]:
        """Process a document file and return its text content as chunks."""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_extension = os.path.splitext(file_path)[1].lower()
        if file_extension not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {file_extension}")

        if file_extension == '.pdf':
            return self._process_pdf(file_path)
        else:  # PPT or PPTX
            return self._process_ppt(file_path)

    def _process_pdf(self, file_path: str) -> List[str]:
        """Extract text from PDF file and return as chunks."""
        chunks = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text.strip():
                    chunks.append(text)
        return chunks

    def _process_ppt(self, file_path: str) -> List[str]:
        """Extract text from PowerPoint file and return as chunks."""
        chunks = []
        prs = Presentation(file_path)
        
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            
            if slide_text:
                chunks.append("\n".join(slide_text))
        
        return chunks 